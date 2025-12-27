"""
Khuta App - PowerPoint Presentation Generator
Creates a comprehensive presentation for the ADHD Assessment App
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
PRIMARY_BLUE = RGBColor(66, 153, 225)
DARK_BLUE = RGBColor(45, 55, 72)
LIGHT_GRAY = RGBColor(247, 250, 252)
GREEN = RGBColor(72, 187, 120)
YELLOW = RGBColor(236, 201, 75)
ORANGE = RGBColor(237, 137, 54)
RED = RGBColor(245, 101, 101)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = PRIMARY_BLUE
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Left bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_BLUE
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    return slide

def add_content_slide(title, content_items, rtl=False):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BLUE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if rtl:
        p.alignment = PP_ALIGN.RIGHT

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(12)
        if rtl:
            p.alignment = PP_ALIGN.RIGHT

    return slide

def add_two_column_slide(title, left_items, right_items, left_title="", right_title="", rtl=False):
    """Add a two-column slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BLUE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if rtl:
        p.alignment = PP_ALIGN.RIGHT

    # Left column title
    if left_title:
        left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5))
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(8)

    # Right column title
    if right_title:
        right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(0.5))
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(7), Inches(2.2), Inches(5.5), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(8)

    return slide

def add_table_slide(title, headers, rows, rtl=False):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BLUE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if rtl:
        p.alignment = PP_ALIGN.RIGHT

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = Inches(12)
    table_height = Inches(0.5 * num_rows)

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.667), Inches(1.8),
        table_width, table_height
    ).table

    # Set column widths
    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)

    # Header row
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_BLUE
            p.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return slide

def add_architecture_slide():
    """Add architecture diagram slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BLUE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture - هيكل النظام"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Architecture boxes
    layers = [
        ("Presentation Layer\nواجهة المستخدم", Inches(1), Inches(1.8), PRIMARY_BLUE),
        ("State Management\nإدارة الحالة (BLoC)", Inches(1), Inches(3.2), RGBColor(129, 140, 248)),
        ("Service Layer\nطبقة الخدمات", Inches(1), Inches(4.6), RGBColor(52, 211, 153)),
        ("Data Layer\nطبقة البيانات", Inches(1), Inches(6), RGBColor(251, 191, 36)),
    ]

    for text, left, top, color in layers:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5), Inches(1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Right side - components
    components = [
        ("Screens & Widgets", Inches(7), Inches(1.8)),
        ("AuthCubit | ChildCubit | AssessmentCubit", Inches(7), Inches(3.2)),
        ("ScoringService | AIService | ErrorHandler", Inches(7), Inches(4.6)),
        ("Firebase | Firestore | Gemini AI", Inches(7), Inches(6)),
    ]

    for text, left, top in components:
        box = slide.shapes.add_textbox(left, top, Inches(5.5), Inches(1))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.LEFT

    # Arrows
    for i in range(3):
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW, Inches(3.25), Inches(2.9 + i * 1.4), Inches(0.5), Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DARK_BLUE
        arrow.line.fill.background()

    return slide

def add_flow_slide():
    """Add user flow slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BLUE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "User Flow - مسار المستخدم"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Flow steps
    steps = [
        ("1. فتح التطبيق", "Splash Screen"),
        ("2. التعريف", "Onboarding"),
        ("3. تسجيل الدخول", "Login/Register"),
        ("4. الشاشة الرئيسية", "Home Screen"),
        ("5. إضافة طفل", "Add Child"),
        ("6. بدء التقييم", "Start Assessment"),
        ("7. الإجابة", "27 Questions"),
        ("8. النتائج", "Results + AI"),
    ]

    x_start = Inches(0.5)
    y = Inches(2.5)
    box_width = Inches(1.4)
    box_height = Inches(1.5)
    gap = Inches(0.15)

    for i, (ar_text, en_text) in enumerate(steps):
        x = x_start + i * (box_width + gap)

        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, box_height
        )
        color = PRIMARY_BLUE if i % 2 == 0 else RGBColor(99, 179, 237)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # Text
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ar_text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = en_text
        p2.font.size = Pt(11)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER

        # Arrow
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, x + box_width, y + Inches(0.6), gap, Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_BLUE
            arrow.line.fill.background()

    return slide

def add_score_interpretation_slide():
    """Add score interpretation slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BLUE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Score Interpretation - تفسير الدرجات"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Score ranges
    scores = [
        ("< 45", "Average\nمتوسط", GREEN, "Low Concern\nقلق منخفض"),
        ("45 - 55", "Slightly Elevated\nمرتفع قليلاً", YELLOW, "Monitor\nيحتاج متابعة"),
        ("55 - 65", "Elevated\nمرتفع", ORANGE, "Concern\nيحتاج اهتمام"),
        ("> 65", "Very Elevated\nمرتفع جداً", RED, "High Concern\nيحتاج تدخل"),
    ]

    y = Inches(1.8)
    for score_range, label, color, desc in scores:
        # Color bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(1.5), Inches(1)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        tf = bar.text_frame
        p = tf.paragraphs[0]
        p.text = score_range
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(Inches(2.3), y, Inches(4), Inches(1))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        # Description
        desc_box = slide.shapes.add_textbox(Inches(7), y, Inches(5.5), Inches(1))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = color

        y += Inches(1.3)

    return slide


# ==================== CREATE SLIDES ====================

# Slide 1: Title
add_title_slide(
    "تطبيق خطى",
    "Khuta - ADHD Assessment App\nتقييم اضطراب فرط الحركة وتشتت الانتباه"
)

# Slide 2: Overview
add_section_slide("نظرة عامة - Overview")

# Slide 3: About the App
add_content_slide(
    "About Khuta - عن تطبيق خطى",
    [
        "تطبيق موبايل لتقييم اضطراب ADHD باستخدام مقياس كونرز",
        "Mobile app for ADHD assessment using Conners' Rating Scale",
        "يدعم اللغتين العربية والإنجليزية",
        "توصيات ذكية باستخدام الذكاء الاصطناعي (Gemini AI)",
        "تقارير PDF قابلة للمشاركة",
        "يعمل بدون إنترنت مع مزامنة تلقائية",
    ]
)

# Slide 4: Features
add_two_column_slide(
    "Key Features - الميزات الرئيسية",
    [
        "تسجيل حساب آمن",
        "إضافة ملفات الأطفال",
        "تقييم الوالدين (27 سؤال)",
        "تقييم المعلم (27 سؤال)",
        "حساب T-Score",
    ],
    [
        "توصيات AI مخصصة",
        "سجل التقييمات السابقة",
        "تقارير PDF",
        "الوضع الليلي",
        "دعم وضع عدم الاتصال",
    ],
    "Authentication & Assessment", "Reports & Settings"
)

# Slide 5: Technology
add_section_slide("التقنيات المستخدمة - Technology Stack")

# Slide 6: Tech Stack Table
add_table_slide(
    "Technology Stack - التقنيات",
    ["Component", "Technology", "المكون"],
    [
        ["Frontend", "Flutter (Dart)", "الواجهة الأمامية"],
        ["State Management", "BLoC / Cubit", "إدارة الحالة"],
        ["Backend", "Firebase", "الخدمات الخلفية"],
        ["Database", "Cloud Firestore", "قاعدة البيانات"],
        ["Authentication", "Firebase Auth", "المصادقة"],
        ["AI", "Google Gemini 2.0", "الذكاء الاصطناعي"],
        ["Reports", "PDF Generation", "التقارير"],
    ]
)

# Slide 7: Architecture Section
add_section_slide("هيكل النظام - System Architecture")

# Slide 8: Architecture Diagram
add_architecture_slide()

# Slide 9: User Flow Section
add_section_slide("مسار المستخدم - User Flow")

# Slide 10: User Flow Diagram
add_flow_slide()

# Slide 11: Assessment Section
add_section_slide("عملية التقييم - Assessment Process")

# Slide 12: Assessment Process
add_content_slide(
    "Assessment Process - عملية التقييم",
    [
        "اختيار نوع التقييم (والدين / معلم)",
        "الإجابة على 27 سؤال من مقياس كونرز",
        "خيارات الإجابة: (0) أبداً - (1) قليلاً - (2) كثيراً - (3) كثيراً جداً",
        "حساب الدرجة الخام من مجموع الإجابات",
        "تحويل الدرجة إلى T-Score حسب العمر والجنس",
        "الحصول على توصيات مخصصة من الذكاء الاصطناعي",
        "حفظ النتائج وإمكانية تصدير تقرير PDF",
    ]
)

# Slide 13: Score Interpretation
add_score_interpretation_slide()

# Slide 14: Database Section
add_section_slide("قاعدة البيانات - Database Schema")

# Slide 15: Database Structure
add_content_slide(
    "Database Structure - هيكل البيانات",
    [
        "Users Collection: بيانات المستخدمين (البريد الإلكتروني، الاسم)",
        "Children Collection: بيانات الأطفال (الاسم، العمر، الجنس)",
        "TestResults Collection: نتائج التقييمات (الدرجة، التوصيات)",
        "Cloud Firestore مع دعم وضع عدم الاتصال",
        "قواعد أمان Firestore لحماية البيانات",
        "المستخدم يمكنه الوصول فقط لبياناته الخاصة",
    ]
)

# Slide 16: Security Section
add_section_slide("الأمان - Security")

# Slide 17: Security Features
add_two_column_slide(
    "Security Features - ميزات الأمان",
    [
        "Firebase Authentication",
        "Email Verification",
        "Password Reset",
        "Secure Session Management",
    ],
    [
        "Firebase App Check",
        "Firestore Security Rules",
        "Data Encryption",
        "Offline Data Protection",
    ],
    "Authentication", "Data Protection"
)

# Slide 18: AI Section
add_section_slide("الذكاء الاصطناعي - AI Recommendations")

# Slide 19: AI Features
add_content_slide(
    "AI-Powered Recommendations - التوصيات الذكية",
    [
        "استخدام Google Gemini 2.0 Flash للتوصيات",
        "تحليل إجابات التقييم لفهم نمط السلوك",
        "توصيات مخصصة حسب درجة التقييم",
        "دعم اللغتين العربية والإنجليزية",
        "توصيات احتياطية في حالة فشل الاتصال",
        "إعادة المحاولة تلقائياً (Retry Logic)",
    ]
)

# Slide 20: Screens Section
add_section_slide("شاشات التطبيق - App Screens")

# Slide 21: Screen List
add_two_column_slide(
    "App Screens - شاشات التطبيق",
    [
        "Splash Screen - شاشة البداية",
        "Onboarding - شاشات التعريف (3)",
        "Login - تسجيل الدخول",
        "Register - إنشاء حساب",
        "Email Verification - التحقق من البريد",
    ],
    [
        "Home - الشاشة الرئيسية",
        "Add Child - إضافة طفل",
        "Child Details - تفاصيل الطفل",
        "Assessment - التقييم (27 سؤال)",
        "Results - النتائج والتوصيات",
        "Settings - الإعدادات",
    ],
    "Authentication Screens", "Main Screens"
)

# Slide 22: Offline Support
add_section_slide("دعم وضع عدم الاتصال - Offline Support")

# Slide 23: Offline Features
add_content_slide(
    "Offline Support - العمل بدون إنترنت",
    [
        "Firestore Persistence مع تخزين محلي غير محدود",
        "عرض البيانات المحفوظة عند انقطاع الاتصال",
        "Offline Queue لحفظ العمليات المعلقة",
        "مزامنة تلقائية عند عودة الاتصال",
        "Offline Banner لإظهار حالة الاتصال",
        "تجربة مستخدم سلسة في جميع الأحوال",
    ]
)

# Slide 24: Testing Section
add_section_slide("الاختبارات - Testing")

# Slide 25: Testing
add_content_slide(
    "Testing Strategy - استراتيجية الاختبارات",
    [
        "Unit Tests: اختبار الخدمات والـ Cubits",
        "Widget Tests: اختبار مكونات الواجهة",
        "Integration Tests: اختبار تدفق العمليات الكاملة",
        "Mock Objects: استخدام Mockito للاختبارات المعزولة",
        "Test Coverage: تغطية شاملة للكود الأساسي",
    ]
)

# Slide 26: Summary
add_section_slide("الملخص - Summary")

# Slide 27: Project Summary
add_content_slide(
    "Project Summary - ملخص المشروع",
    [
        "تطبيق Flutter متكامل لتقييم ADHD",
        "استخدام BLoC/Cubit لإدارة الحالة",
        "Firebase للمصادقة وتخزين البيانات",
        "Gemini AI للتوصيات الذكية",
        "دعم اللغتين العربية والإنجليزية",
        "تصميم حديث مع Dark/Light Mode",
        "يعمل بدون إنترنت مع مزامنة تلقائية",
    ]
)

# Slide 28: Thank You
add_title_slide(
    "شكراً لكم",
    "Thank You\n\nKhuta - ADHD Assessment App"
)

# Save presentation
output_path = "Khuta_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
